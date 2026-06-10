"""
Generate teaching PPT for Jyotish-Prasana Chapter 1.
Structured into 6 pedagogical modules.
Source text is curated into presentation-friendly format (large fonts, key points).
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parent.parent
INPUT_FILE = PROJECT_ROOT / "output" / "Jyotish-Prasana.txt"
OUTPUT_FILE = PROJECT_ROOT / "output" / "Jyotish-Prasana_Ch1_v3.pptx"

MAX_TABLE_COLS = 8
BODY_FONT = Pt(18)
SMALL_FONT = Pt(14)
TITLE_FONT = Pt(30)
TABLE_FONT = Pt(10)
ACCENT_COLOR = RGBColor(0x8B, 0x00, 0x00)  # dark red

# ============================================================
# Slide content — hand-curated from source text
# ============================================================

# Each slide: (title, body_text, [tables_from_source])

INTRO_SLIDES = [
    ("什么是 Muhurtha（择时）？", (
        "Muhurtha 指任何一个短暂的时间段，无论吉凶。\n\n"
        "时间有两种属性：\n"
        "  • 无属性（Nirguna）—— 无限、永恒，宇宙在其中创生与消融\n"
        "  • 有属性（Saguna）—— 可测量、变动，我们日常体验的时间\n\n"
        "有属性时间像车轮一样在太阳的控制下转动。\n"
        "因此，一切时间划分都依赖太阳在黄道上的公转。"
    )),
    ("时间划分的五个层级", (
        "时间单位从大到小，最重要的五个：\n\n"
        "    年 → 月 → 日 → 上升星座（Lagna）→ 时（Muhurta）\n\n"
        "太阳完成一次公转 = 1 年。\n"
        "由此计算出月、日、Ghati（24分钟）、Vighati（24秒）等。\n\n"
        "所有活动都离不开 Muhurta-shuddhi（时之纯净）。\n"
        "择时的宗旨：多优点，少瑕疵。\n"
        "连造物主也无法找到一个绝对完美的时间。"
    )),
]

TIME_SLIDES = [
    ("太阴年（Lunar Years）", (
        "共有 60 个太阴年，构成一个完整周期。\n"
        "Lagadha（吠陀占星学作者）引入 Yuga = 5 年周期：\n"
        "  Samvatsara → Anuvatsara → Parivatsara → Idavatsara → Idavatsara\n\n"
        "60 年 = 12 个五年周期，分别由 12 位神祇掌管。"
    )),
    ("季节（Ritu）与太阳月", (
        "一年分为 6 个季节，每个季节对应 2 个太阳月：\n\n"
        "  春季：双鱼座 + 白羊座（Chaitra + Vaisakha）\n"
        "  夏季：金牛座 + 双子座（Jyeshtta + Ashada）\n"
        "  雨季：巨蟹座 + 狮子座（Sravana + Bhadrapada）\n"
        "  秋季：处女座 + 天秤座（Aswiyuja + Kartika）\n"
        "  晚秋：天蝎座 + 射手座（Margasira + Paushya）\n"
        "  冬季：摩羯座 + 水瓶座（Magha + Phalguna）\n\n"
        "前 3 季属神性季节，后 3 季与祖灵相关。"
    ), []),
    ("阿衍那（Ayana）—— 太阳的南北路径", (
        "Ayana = 太阳在赤道南北的半年路径。\n\n"
        "Uttarayana（夏至 → 冬至，12/21 至 6/21）：\n"
        "  涵盖冬季 + 春季 + 夏季\n"
        "  适宜：圣线仪式、婚礼、乔迁、曼陀罗灌顶等\n\n"
        "Dakshinayana（6/21 → 12/21）：\n"
        "  涵盖雨季 + 秋季 + 晚秋\n\n"
        "Uttarayana 开始后 6h24m 内 = Punyakala（吉祥时刻）\n"
        "此期间必须绝对避免任何活动。"
    ), []),
    ("月份的四类划分", (
        "1. 太阳月：太阳进入一个星座，约 30天10小时\n"
        "2. 太阴月（朔望月）：新月到新月，约 29天12小时\n"
        "3. Savana 月：固定 30 天，每天 60 Ghati\n"
        "4. 恒星月：月亮通过 27 宿，约 27天8小时\n\n"
        "秽月（Malamasas）：\n"
        "  • 闰月（Adhika masa）：一个月内无太阳进入\n"
        "  • 缺月（Kshaya masa）：一个月内两次太阳进入\n"
        "  秽月期间禁止一切吉祥活动。"
    ), []),
    ("半月（Paksha）—— 白与黑", (
        "Sukla Paksha（白半月）：新月 → 满月，15 天\n"
        "Krishna Paksha（黑半月）：满月 → 新月，15 天\n\n"
        "若一个半月内发生 2 次 Tithi 缺失，\n"
        "可能只有 13 天 → 称为 Raura Kalayoga（地狱时段）\n"
        "→ 禁止吉祥活动。"
    ), []),
]

PANCHANGA_SLIDES = [
    ("Panchanga（五支）概述", (
        "Panchanga = 五个组成部分，是择时评估的核心框架：\n\n"
        "  1. Tithi（太阴日）—— 月亮与太阳的角度差\n"
        "  2. Vara（星期）—— 七日周期\n"
        "  3. Nakshatra（星宿）—— 27 个太阴宿\n"
        "  4. Yoga（瑜伽）—— 日月经度组合关系\n"
        "  5. Karana（迦罗那）—— 半个太阴日\n\n"
        "每一项都有：计算方法 → 吉凶分类 → 禁忌时段 → 例外与补救"
    ), []),
    ("Tithi（太阴日）的计算方法", (
        "Tithi 的定义：月亮与太阳的经度差。\n\n"
        "月亮每天比太阳多走 12°（月亮 13°/天，太阳 1°/天）\n"
        "30 天 × 12° = 360° = 一个朔望月\n\n"
        "计算步骤：\n"
        "  1. 求月亮从太阳处移动的度数\n"
        "  2. 除以 12°（720'）\n"
        "  3. 商 + 1 = Tithi 序号（≤15 = 白半月，>15 = 黑半月）\n"
        "  4. 用比例法计算剩余部分的结束时间"
    ), []),
    ("Tithi 的五种分类", (
        "每个 Tithi 归属于五种类型之一：\n\n"
        "  Nanda（喜悦）：第 1、6、11 日 → 穿新衣、学音乐舞蹈\n"
        "  Bhadra（吉祥）：第 2、7、12 日 → 圣线礼、婚礼、佩戴饰品\n"
        "  Jaya（胜利）：第 3、8、13 日 → 入伍、建房、乔迁、行商\n"
        "  Rikta（空）：第 4、9、14 日 → 手术、逮捕（凶事专用）\n"
        "  Poorna（圆满）：第 5、10、15 日 → 婚礼、祭祀（新月除外）\n\n"
        "口诀：Nanda-Bhadra-Jaya-Rikta-Poorna 五字诀。"
    ), []),
    ("缺陷 Tithi（Chidra Tithis）", (
        "第 4、6、8、9、12、14 日 = 缺陷 Tithi（Paksharandhra Tithis）\n"
        "通常不适用于吉祥活动。\n\n"
        "各 Tithi 应避免的 Ghati 数（1 Ghati = 24 分钟）：\n"
        "  第 4 日：8 Ghati    第 8 日：14 Ghati   第 12 日：10 Ghati\n"
        "  第 6 日：9 Ghati    第 9 日：25 Ghati   第 14 日：5 Ghati\n\n"
        "若 Tithi 实际时长 ≠ 60 Ghati，用比例法重新计算禁忌时段。"
    ), []),
    ("Tithi 禁忌时段（Visha Ghatis）", (
        "每个 Tithi 有其特定的 Visha Ghatis（毒时），应避免使用。\n\n"
        "第 1 日：15   第 6 日：5    第 11 日：3\n"
        "第 2 日：5    第 7 日：4    第 12 日：10\n"
        "第 3 日：8    第 8 日：8    第 13 日：12\n"
        "第 4 日：7    第 9 日：7    第 14 日：?\n"
        "第 5 日：7    第 10 日：10   第 15 日：8\n\n"
        "补救：月亮在 4/5/9/10/11 宫，或命宫主星在四正宫，\n"
        "或受吉星相位照射，可减轻毒时影响。"
    ), []),
    ("Tithi 的选择原则", (
        "一个 Tithi 通常跨越两天的部分时段。\n\n"
        "民事 / 斋戒 / 吉祥活动：以日出时正在进行的 Tithi 为准\n"
        "宗教活动：可能需在 Tithi 开始的前一天庆祝\n"
        "出生 / 死亡：使用即时 Tithi（Tatkalika Tithi）\n\n"
        "缺失 Tithi（Kshaya）和增长 Tithi（Vriddhi）：\n"
        "  均不被接受用于吉祥活动。\n"
        "  补救：择时盘中木星位于四正宫（Kendra）。"
    ), []),
    ("Vara（星期）的吉凶", (
        "吉日：星期一、三、四、五\n"
        "凶日：星期二、六、日（但夜间不算凶）\n\n"
        "星期毒时（Vara Visha Ghatis）：\n"
        "  日 20、月 2、火 12、水 10、木 7、金 5、土 25\n"
        "  各持续 4 Ghati（1h36m）\n\n"
        "星期凶时（Vara Durmuhurtha）：\n"
        "  一个恒星日 = 30 Muhurtha（每个 2 Ghati = 48 分钟）\n"
        "  每天有特定编号的 Muhurtha 为凶时。"
    ), []),
    ("Durmuhurtha 速查表", (
        "每天 30 个 Muhurtha（白天 15 + 夜晚 15），各星期的凶时 Muhurtha：\n\n"
        "  星期日：第 14 个（Aryama）\n"
        "  星期一：第 9（Brahma）+ 第 12（Asura）\n"
        "  星期二：第 4（Pitrya）+ 夜间第 7（Anala）\n"
        "  星期三：第 8（Abhijit）\n"
        "  星期四：第 6（Vara）+ 第 12（Asura）\n"
        "  星期五：第 4（Pitrya）+ 第 9（Brahma）\n"
        "  星期六：第 1（Rudra）+ 第 2（Sarpa）"
    ), []),
    ("白昼四种凶时（Diurnal Evil Periods）", (
        "四种凶时各有不同来源和争议：\n\n"
        "  Gulika（Kulika）：普遍认为对不吉仪式为凶，吉祥活动无碍\n"
        "  Yamaghanta（Yamakantaka）：仅周四不吉（木星之子）\n"
        "  Kalavela（Ardhayama）：前 48 分钟为凶\n"
        "  Kantaka（Mrityu）：象征死亡\n\n"
        "日出到日落 ÷ 8 = 每部分时长。\n"
        "不同星期凶时所在部分各有不同（见源文表格）。\n\n"
        "重要：若吉星在四正宫，上述凶时均可减轻或消除。\n"
        "日常活动无需考虑以上凶时。"
    ), []),
    ("Nakshatra（星宿）—— 27 个太阴宿", (
        "黄道被等分为 27 个星宿，每个 13°20'（800'）。\n\n"
        "计算方法：将行星经度（角分）÷ 800\n"
        "  商 = 已穿越的星宿数量\n"
        "  余数 = 当前星宿中已穿越的部分\n\n"
        "有些体系包含第 28 宿 Abhijit（摩羯座 9°6' - 10°53'20\"）。\n\n"
        "星宿的结束时刻 = 将尚未经过的部分 ÷ 日行度数 × 24h"
    ), []),
    ("星宿的 14 种分类法（上）", (
        "固定星宿（Dhruva）：Rohini, Uttara, Uttarashada, Uttarabhadra\n"
        "  → 播种、穿新衣、入宅、学音乐\n\n"
        "移动星宿（Chara）：Punarvasu, Swati, Sravana, Dhanishta, Satabhisha\n"
        "  → 驾车、开店、制作金器银器\n\n"
        "光明星宿（Kshipra）：Aswini, Pushya, Hasta, Abhijit\n"
        "  → 商品销售、学习科学、药物治疗\n\n"
        "锐利星宿（Tikshna）：Ardra, Aslesha, Moola, Jyeshtha\n"
        "  → 训练象马、播种、乔迁、开始教育\n\n"
        "凶猛星宿（Ugra）：Bharani, Magha, Purvaphalguni, Purvashada, Purvabhadra\n"
        "  → 邪恶活动及锐利星宿所述活动"
    ), []),
    ("星宿的 14 种分类法（下）", (
        "面朝上星宿（Urdhvamukha）：建房、婚姻、宣誓典礼\n"
        "面朝下星宿（Adhomukha）：挖井、学习占星学\n"
        "斜向星宿（Tiryagmukha）：耕作、驾摩托车、航空/水路旅行\n\n"
        "盲星宿（Andha）：失物向东找回\n"
        "弱视星宿（Mandaksha）：失物向南，极难寻回\n"
        "中视星宿（Madhyaksha）：失物向西\n"
        "善视星宿（Sulochana）：失物向北，无寻回可能\n\n"
        "四足 / 异足 / 等足：用于特定仪式判断。"
    ), []),
    ("Varjyam（弃时）与 Amrita Kalam（甘露时）", (
        "Varjyam = 每个星宿中必须避开的不吉时段。\n"
        "计算方法：取星宿总时长的特定分数（各宿不同）。\n\n"
        "Pushyami 示例：\n"
        "  Varjyam = 星宿总时长的 1/3\n"
        "  → 从 Punarvasu 结束 + Varjyam 偏移 = 上午 10:22 - 12:02\n\n"
        "Amrita Kalam = 约 4 Ghati（1h36m）的吉时。\n"
        "从 Varjyam 起点加减特定 Ghati 数计算。\n\n"
        "补救：月亮在 4/7/10/5/9 宫，或命宫主星在 1/4/5/7/9/10 宫。"
    ), []),
    ("Tarabalam（星宿力）与 Chandrabalam（月亮力）", (
        "Tarabalam：从出生星宿数到当日星宿，除以 9，看余数：\n"
        "  1=出生 2=财富 3=灾祸 4=安稳 5=障碍\n"
        "  6=成就 7=死亡 8=友好 9=至友\n"
        "  → 2、4、6、8、9 吉；1、3、5、7 凶\n\n"
        "Chandrabalam：月亮所在星座是出生星座的第几位：\n"
        "  第 1/3/6/7/10/11 位 → 有月亮力\n"
        "  第 2/5/9 位 → 白半月吉；第 4/8/12 位 → 黑半月吉\n\n"
        "建议：黑半月求 Tarabalam，白半月求 Chandrabalam。"
    ), []),
    ("Yoga（瑜伽）—— 27 种先天瑜伽", (
        "Yoga = 太阳经度 + 月亮经度的组合。\n"
        "先天瑜伽（Naisergika）= 按 13°20' 间隔划分 360° = 27 种。\n\n"
        "1. Vishkambha   8. Dhriti    15. Vajra     22. Sadhya\n"
        "2. Preeti       9. Soola     16. Siddhi    23. Subha\n"
        "3. Ayushman    10. Ganda     17. Vyatipata 24. Sukla\n"
        "4. Sowbhagya   11. Vriddhi   18. Variyan   25. Brahman\n"
        "5. Sobhana     12. Dhruva    19. Parigha   26. Indra\n"
        "6. Atiganda    13. Vyaghata  20. Siva      27. Vaidhriti\n"
        "7. Sukarma     14. Harshana  21. Siddha\n\n"
        "计算方法：(太阳经度 + 月亮经度) ÷ 13°20' → 商 + 1 = Yoga 序号"
    ), []),
    ("即时瑜伽（Tatkalika Yogas）与吉凶组合", (
        "即时瑜伽有 28 种（Ananda → Vardhamana），\n"
        "基于当日星宿即时判断。\n\n"
        "凶险瑜伽由元素组合而成：\n"
        "  Tithi + Vara 形成的凶瑜伽\n"
        "  Tithi + Vara 形成的吉瑜伽\n"
        "  Vara + Nakshatra 形成的吉瑜伽\n"
        "  Vara + Nakshatra 形成的凶瑜伽（Yamaghatta, Dagdha, Utpata）\n"
        "  Tithi + Nakshatra + Vara 三重组合的凶瑜伽\n\n"
        "每种凶瑜伽都有例外和补救措施（详见表）。"
    ), []),
    ("Karana（迦罗那）—— 半个太阴日", (
        "Karana = Tithi 的后半段，每个 Tithi 有两个 Karana。\n"
        "一个朔望月共 11 种 Karana：\n\n"
        "  固定 7 种（每月重复 8 次）：\n"
        "  Bava, Balava, Kaulava, Taitila, Gara, Vanija, Vishti\n\n"
        "  变动 4 种（每月出现 1 次）：\n"
        "  Kimstughna, Sakuni, Chatushpada, Naga\n\n"
        "Bhadra（Vishti Karana 的别称）：\n"
        "  第 8、9 日黑半月的后半个 Karana 特别不吉\n"
        "  → 禁止一切吉祥活动。"
    ), []),
]

PLANET_SLIDES = [
    ("行星在黄道中的关系", (
        "主宰星座（Rulership）：\n"
        "  Su → 狮子    Mo → 巨蟹    Ma → 白羊/天蝎    Me → 双子/处女\n"
        "  Ju → 射手/双鱼    Ve → 金牛/天秤    Sa → 摩羯/水瓶\n\n"
        "擢升（Exaltation）：行星力量最强的位置\n"
        "落陷（Debilitation）：擢升的对宫，力量最弱\n"
        "Moolatrikona：行星的根本三角位\n\n"
        "相位：所有行星均相位其第 7 宫。\n"
        "  额外特殊相位：木星 +5+9、火星 +4+8、土星 +3+10"
    ), []),
    ("十二宫的分类", (
        "Kendras（四正宫）：第 1、4、7、10 宫\n"
        "Trikonas（三方宫）：第 1、5、9 宫\n"
        "Panaparas（续宫）：第 2、5、8、11 宫\n"
        "Apoklimas（果宫）：第 3、6、9、12 宫\n"
        "Upachayas（增长宫）：第 3、6、10、11 宫\n"
        "Trikas（凶宫）：第 6、8、12 宫"
    ), []),
    ("黄道星座的 12 种分类", (
        "1. 移动：白羊/巨蟹/天秤/摩羯\n"
        "2. 固定：金牛/狮子/天蝎/水瓶\n"
        "3. 双体：双子/处女/射手/双鱼\n"
        "4. 阳性：白羊/双子/狮子/天秤/射手/水瓶\n"
        "5. 阴性：金牛/巨蟹/处女/天蝎/摩羯/双鱼\n"
        "6. 尾部升起（Prishttodaya）\n"
        "7. 头部升起（Seershodaya）\n"
        "8. 双部升起（Ubhayodaya）\n"
        "9-12. 四姓分类：婆罗门/刹帝利/吠舍/首陀罗"
    ), []),
    ("太阳入宫（Solar Ingress）", (
        "太阳从一个星座进入另一个 = Sankranti，约 30 天一次。\n\n"
        "特殊入宫类型：\n"
        "  Shadasiti：太阳入双子/处女/射手/双鱼\n"
        "  Vishu：太阳入白羊/天秤（二分点）\n"
        "  Vishnupada：太阳入金牛/狮子/天蝎/水瓶\n"
        "  Yamayana / Saumyayana：南行/北行\n\n"
        "Punyakala（吉祥时段）：\n"
        "  午夜前入宫 → Punyakala 属前一天后期\n"
        "  午夜后入宫 → Punyakala 属第二天前期"
    ), []),
]

SPECIAL_SLIDES = [
    ("木星偕日升系统与燃烧期", (
        "木星 12 年周期（Jovian Years）影响择时质量。\n\n"
        "木星与金星有幼年、老年和燃烧期（Moudya）：\n"
        "  燃烧期 = 行星距离太阳过近、无法被观测的阶段\n"
        "  太阳 - 行星经度差 < 特定度数 → 燃烧\n\n"
        "燃烧期禁止：十六大布施、苏摩祭祀、四月斋戒、\n"
        "学成归家（Samavartana）、出家受戒等。\n\n"
        "Gurumoudyam（木星燃烧）：\n"
        "  约持续 1-2 月不等，期间不可举行婚礼、圣线礼。\n"
        "  朝圣（Gaya/Godavari）不受燃烧限制。"
    ), []),
    ("狮子座木星之过（Simha-Guru Dosha）", (
        "木星行经狮子座 Magha 四足 + Pubba 第一足期间：\n"
        "  → 全印度不适宜圣线仪式和婚姻。\n\n"
        "木星在 Pubba 第一足达到 Vargottamamsa 时：\n"
        "  → 任何地区不得举行婚礼（可能导致夫妻死亡）。\n\n"
        "Pubba 第二/三/四足 + Uttara 第一足：\n"
        "  → 在 Andhra/Karnataka/Tamilnadu/Maharashtra 等地不吉\n"
        "  → 在 Haridwar/Almorah/Nainital 等地例外。"
    ), []),
    ("日月食（Eclipses）的禁忌", (
        "日月食是择时中最重要的否决项之一。\n\n"
        "食当日：禁止一切吉祥仪式。\n\n"
        "食后星宿禁忌期：\n"
        "  全食 → 6 个月内避免使用该星宿\n"
        "  半食 → 3 个月\n"
        "  四分之一食 → 1 个月\n\n"
        "食前后禁忌：\n"
        "  偏食 → 前后 3 天排除\n"
        "  全食 → 前后 7 天排除\n"
        "  食始（Grasthastha）→ 前 3 天排除\n"
        "  食终（Grasthodaya）→ 后 3 天排除"
    ), []),
    ("霍拉什塔卡（Holashtaka）", (
        "从 Phalguna 月白半月第 8 Tithi 至 Holi 节的 8 天。\n\n"
        "影响区域（主要在 Punjab/Himachal 地区）：\n"
        "  Hoshiarpur, Gurdaspur, Mandee, Kapurthala,\n"
        "  Lahore, Amritsar, Ferojpur, Simla, Ludhiana, Ajmer 等。\n\n"
        "这些区域在此期间应避免婚姻等吉祥仪式。\n\n"
        "这是择时中按地理区域差异化处理的典型案例。"
    ), []),
]

INTEGRATION_SLIDES = [
    ("Panchaka-rahitam（五数除凶）", (
        "将 Tithi + Vara + Nakshatra + Muhurtha Lagna 数值化\n"
        "在五个位置进行加法和除以 9 的判断。\n\n"
        "五个位置的余数含义：\n"
        "  位置 I   → Raga Panchaka（贪五数）\n"
        "  位置 II  → Nripa Panchaka（王五数）\n"
        "  位置 III → Chora Panchaka（贼五数）\n"
        "  位置 IV  → Mrityu Panchaka（死五数）\n"
        "  位置 V   → （不吉）\n\n"
        "简化法：总和 ÷ 9，余数 1/2/4/6/8 = 不吉，3/5/7/9 = 吉。\n"
        "死五数在任何情况下都应避免。"
    ), []),
    ("Pancheshtika（五砖评估法）", (
        "五颗行星 × 20 分 = 100 分制的择时评分系统。\n\n"
        "评分项：\n"
        "  1. Kaladhipa（时主星）—— 20 分\n"
        "  2. Lagna Lord（命宫主星）—— 20 分\n"
        "  3. Jupiter（木星）—— 20 分\n"
        "  4. Moon（月亮）—— 20 分\n"
        "  5. Sun（太阳）—— 20 分\n\n"
        "评分标准：\n"
        "  5 颗均佳 = 100%（Pancheshtika）\n"
        "  4 颗佳（不含太阳）= 80%\n"
        "  3 颗佳（时主+命主+木星）= 60%（Triashtika，最低可接受）\n"
        "  吉星应在擢升/本垣/友垣或四正宫/四角宫\n"
        "  凶星应在 3/6/11 宫（但命宫主星为凶星时不应在第 6 宫）"
    ), []),
    ("Pushkaramsa 择时——精确到分钟", (
        "择时上升的强度取决于其 Navamsa（九分盘）。\n"
        "吉星 Navamsa：双子、处女、天秤、金牛。\n\n"
        "Pushkara 度数（最优窗口，约 3-4 分钟）：\n"
        "  火象（白羊/狮子/射手）：21°\n"
        "  土象（金牛/处女/摩羯）：14°\n"
        "  风象（双子/天秤/水瓶）：24°\n"
        "  水象（巨蟹/天蝎/双鱼）：7°\n\n"
        "案例：白羊上升 = 上午 6:40-8:25（105 分钟/30°）\n"
        "  Pushkara 21° = 6:40 + 73.5 分钟 = 7:51-7:54（吉时窗口）"
    ), []),
    ("Abhijit Lagna（阿毗吉特上升）", (
        "日出上升起算的第 4 个上升。\n\n"
        "优势：能消除多种凶兆（Utpata、Vishti 等）。\n"
        "可用于婚礼等吉祥仪式。\n\n"
        "但有三条重要限制：\n"
        "  1. 不适用于圣线仪式（Upanayana）——\n"
        "     第 10 宫被太阳占，产生 Sphujita Dosha\n"
        "  2. 星期三禁用\n"
        "  3. 向南旅行禁用\n\n"
        "婚礼需在上午 11:55 之前（Poorvahna），\n"
        "即使 Abhijit 持续到下午 2:00。"
    ), []),
    ("Godhuli Lagna（戈杜利上升）", (
        "日出上升的第 7 个上升，发生在日落时分。\n"
        "得名于牛群归家、蹄扬尘土的景象。\n\n"
        "高度受推崇（Rama Daivajna 在 Muhurtha Chintamani 中赞扬），\n"
        "但有限制：\n\n"
        "  适用：首陀罗、自由恋爱婚姻、跨种姓婚姻、Orissa 地区\n"
        "  不适用：婆罗门婚礼（绝对不可！）\n\n"
        "禁忌：月亮在 6 或 8 宫，或 6/8 宫被火/水/木/金占据\n"
        "星期四日落前、星期六日落后不吉\n"
        "仅当月亮在 2/3/11 宫时才为吉。"
    ), []),
    ("Pushkara Yoga（普什卡拉瑜伽）", (
        "两种类型：\n\n"
        "Dwipushkara Yoga（双倍果瑜伽）—— 27 种变体：\n"
        "  第一组：星期六/日/二\n"
        "  第二组：Bhadra Tithi（第 2/7/12 日）\n"
        "  第三组：双足星宿（Dhanishta/Mrigasira/Chitta）\n"
        "  任意两组的组合即形成 Dwipushkara\n\n"
        "Tripushkara Yoga（三倍果瑜伽）—— 54 种变体：\n"
        "  第三组替换为三足星宿（Visakha/Uttara/Purvabhadra 等）\n\n"
        "注意：吉凶均加倍/三倍！\n"
        "死亡可能接连发生两次或三次。\n"
        "在财务事项中（利息业务、银行存款）极为受推崇。"
    ), []),
    ("Pradosha（普拉多沙）的禁忌", (
        "Pradosha = 特定 Tithi 中的黄昏时段。\n\n"
        "Chaturdashi（第 14 日）前 3 小时\n"
        "Saptami（第 7 日）前 4.5 小时\n"
        "Trayodashi（第 13 日）前 6 小时\n\n"
        "从日落起持续 1h12m - 3h24m 的时段 = Pradosha\n"
        "所有吉祥活动应避免。\n"
        "夜间誓戒（Vrata）不受此限。"
    ), []),
]

# All modules in order
ALL_MODULES = [
    ("intro", "绪论：什么是择时？", INTRO_SLIDES, (
        "择时的本质：在所有可能的时间中，找到一个'多优点、少瑕疵'的窗口。\n"
        "这需要系统化的评估方法 —— 这正是第一章要构建的知识体系。"
    )),
    ("time", "第一部分：时间划分体系", TIME_SLIDES, (
        "时间是有层级结构的。择时从最粗粒度（年）向最细粒度（Muhurta）逐级筛选。\n"
        "理解这个层级，才能理解后续 Panchanga 各要素在哪个粒度上起作用。"
    )),
    ("panchanga", "第二部分：五支（Panchanga）评估体系", PANCHANGA_SLIDES, (
        "Panchanga 是择时判断的五项基本功。\n"
        "每项要素都遵循：计算方法 → 吉凶分类 → 禁忌时段 → 例外与补救 的四步结构。"
    )),
    ("planets", "第三部分：行星因素与宫位", PLANET_SLIDES, (
        "Panchanga 评估的是'时间'质量，行星因素评估的是'空间'状态（天空格局）。\n"
        "二者结合才能做出完整判断。"
    )),
    ("special", "第四部分：特殊天文事件", SPECIAL_SLIDES, (
        "特殊天文事件是择时中的'否决项'—— 一旦触发，即使 Panchanga 再好也不能用。\n"
        "优先检查日月食、燃烧期、区域禁忌。"
    )),
    ("integration", "第五部分：综合评估与高级择时", INTEGRATION_SLIDES, (
        "从定性到定量：Panchanga 是定性判断 → Pancheshtika 是 100 分制评分\n"
        "→ Pushkaramsa 是精确到分钟的最优窗口。"
    )),
]


def clean_source_table(raw_text: str, table_id: str) -> list[list[str]]:
    """Extract a specific table from the source text."""
    text = raw_text
    # Find table by context
    lines = text.split("\n")
    tables = []
    current_table = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            if not re.match(r"^\|[\s\-:|]+\|$", stripped):
                current_table.append([c.strip() for c in stripped[1:-1].split("|")])
            in_table = True
        else:
            if in_table and len(current_table) >= 2:
                tables.append(current_table)
            current_table = []
            in_table = False

    if in_table and len(current_table) >= 2:
        tables.append(current_table)

    # Return the most relevant table (or empty)
    if tables:
        return tables[0] if len(tables[0]) <= 20 else []
    return []


def add_content_slide(prs, title: str, body: str):
    """Add a slide with title and body text. No module tags."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    # Set title font
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        p.font.size = TITLE_FONT

    body_frame = slide.shapes.placeholders[1].text_frame
    body_frame.clear()

    first = True
    for line in body.strip().split("\n"):
        if first:
            p = body_frame.paragraphs[0]
            first = False
        else:
            p = body_frame.add_paragraph()
        p.text = line
        p.font.size = SMALL_FONT if line.startswith("  ") else BODY_FONT
        p.space_after = Pt(6)


def add_module_divider(prs, title: str, overview: str, takeaway: str):
    """Module overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(28)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    p = body.paragraphs[0]
    p.text = overview
    p.font.size = BODY_FONT
    p.space_after = Pt(18)

    p2 = body.add_paragraph()
    p2.text = f"▸ {takeaway}"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_COLOR


def read_chapter_text() -> str:
    with open(INPUT_FILE, "r", encoding="utf-8") as f:
        text = f.read()
    lines = text.split("\n")
    start = end = None
    for i, l in enumerate(lines):
        if "**第一章 绪论**" in l:
            start = i + 1
        if start and l.startswith("**第二章**"):
            end = i
            break
    return "\n".join(lines[start:end]) if start else ""


def main():
    print(f"Reading: {INPUT_FILE}")
    source_text = read_chapter_text()
    print(f"  Chapter 1: {len(source_text)} chars")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Title slide ----
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = "第一章  绪论"
    ts.shapes.placeholders[1].text = (
        "Muhurtha Sindhu（择时宝鉴）\n"
        "作者：Iranganti Rangacharya\n\n"
        "择时占星学的完整知识体系"
    )

    # ---- Architecture overview ----
    arch_text = (
        "本章系统讲授择时占星学的完整评判体系，分五个层次：\n\n"
        "一、时间划分体系\n"
        "    年 → 季 → 月 → 半月 → 太阴日\n"
        "    ↓\n"
        "二、五支（Panchanga）评估\n"
        "    Tithi / Vara / Nakshatra / Yoga / Karana\n"
        "    每一项：计算方法 → 吉凶分类 → 禁忌时段 → 例外与补救\n"
        "    ↓\n"
        "三、行星因素与宫位\n"
        "    行星强弱 + 宫位吉凶 + 星座分类\n"
        "    ↓\n"
        "四、特殊天文事件（否决项）\n"
        "    日月食 / 木星燃烧 / 区域禁忌\n"
        "    ↓\n"
        "五、综合评估与高级择时\n"
        "    定量评分 → 精确到分钟的最优窗口\n\n"
        "核心原则：从粗到细、从定性到定量、从一般规则到特殊例外。"
    )
    add_content_slide(prs, "本章知识架构", arch_text)

    # ---- Module slides ----
    total = 3  # title + architecture + final summary
    for mod_id, mod_title, slides, takeaway in ALL_MODULES:
        add_module_divider(prs, mod_title,
                           f"本部分讲解：{mod_title.split('：')[1] if '：' in mod_title else mod_title}",
                           takeaway)
        total += 1
        for entry in slides:
            if len(entry) == 3:
                title, body, _ = entry
            else:
                title, body = entry
            add_content_slide(prs, title, body)
            total += 1

    # ---- Final summary: complete flow ----
    summary_text = (
        "完整的择时判断流程（第一章总结）：\n\n"
        "第一步：检查否决项\n"
        "  → 日月食？秽月？Holashtaka？燃烧期？\n"
        "  → 有则直接排除，无需继续评估\n\n"
        "第二步：Panchanga 五支逐项评估\n"
        "  → Tithi：类型（Nanda/Bhadra/Jaya/Rikta/Poorna）+ Visha Ghatis\n"
        "  → Vara：吉凶日 + Durmuhurtha + 四种白昼凶时\n"
        "  → Nakshatra：分类 + Varjyam + Tarabala/Chandrabala\n"
        "  → Yoga：先天 27 种 + 即时 28 种，检查凶险组合\n"
        "  → Karana：11 种类型 + Bhadra 禁忌\n\n"
        "第三步：行星因素检查\n"
        "  → Kaladhipa + Lagna Lord + Jupiter + Moon + Sun 的位置\n\n"
        "第四步：综合量化\n"
        "  → Panchaka-rahitam（五数排除法）\n"
        "  → Pancheshtika（100分制，至少 60 分）\n"
        "  → Pushkaramsa（精确到 3-4 分钟的最优窗口）\n\n"
        "最终原则：多优点，少瑕疵。绝对完美的时间不存在。"
    )
    add_content_slide(prs, "第一章总结：择时判断完整流程", summary_text)
    total += 1

    print(f"  Created {len(prs.slides)} slides")

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
