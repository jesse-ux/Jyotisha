"""
Generate teaching PPT for Zoller's Medieval Astrology Foundation Course.
Starts with the Orientation lecture (pages 4-24 of the PDF).

Structure follows Zoller's own logical progression:
1. Why Orientation? — The instructor's perspective and the value of tradition
2. The Hermetic Context — Where Medieval Astrology fits in Western esotericism
3. Medieval vs Modern — Key differences in approach
4. Source Texts — The canon of Medieval Astrology
5. Fate & Free Will — The philosophical foundation
6. Course Overview — What you will learn
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_FILE = PROJECT_ROOT / "output" / "Zoller-Medieval-Astrology_Orientation.pptx"

BODY_FONT = Pt(18)
SMALL_FONT = Pt(14)
TITLE_FONT = Pt(30)
ACCENT_COLOR = RGBColor(0x8B, 0x00, 0x00)

# ============================================================
# Orientation Slides — hand-curated from the course content
# ============================================================

ORIENTATION_SLIDES = [
    # Module 1: Why Orientation?
    ("导论课的必要性", (
        "Robert Zoller 在课程开始前提出三个理由，\n"
        "解释为什么需要一堂导论课：\n\n"
        "a. 具体技法必须在其正确的语境中才能被理解和正确运用\n\n"
        "b. 你需要知道讲师的立场 —— 本课程的观点不是现代的，\n"
        "   对许多人来说可能相当陌生\n\n"
        "c. 这是一门传统西方中世纪占星学课程 ——\n"
        "   Zoller 只教授来自已知来源的传统教义，\n"
        "   不做任何向'当代教条'的妥协\n\n"
        "核心立场：你先了解传统。了解之后，怎么用是你的事。\n"
        "但如果传统中有价值的东西，你切断了与传统的联系就无法接收。"
    )),

    ("中世纪占星学的赫尔墨斯语境", (
        "中世纪占星学不能脱离赫尔墨斯传统（Hermetic Tradition）来理解。\n\n"
        "赫尔墨斯传统包含三门核心艺术：\n"
        "  占星学（Astrology）—— 理解天道\n"
        "  炼金术（Alchemy）—— 转化物质\n"
        "  魔法（Magic）—— 与神灵/智性沟通\n\n"
        "三者的层级关系：\n"
        "  占星学 → 数学 → 对神圣秩序的理性认知\n"
        "  炼金术 → 冶金术 → 对物质世界的掌控\n"
        "  魔法 → 奇迹 → 对超自然力量的沟通\n\n"
        "Zoller 特别引用 Firmicus Maternus（4世纪）：\n"
        "'学习和实践占星学，培养的是宗教和对神圣的崇拜。'"
    )),

    ("高等智慧与世间智慧", (
        "Zoller 区分两种智慧传统：\n\n"
        "高等智慧（Higher Wisdom）：\n"
        "  • 对神圣的直观认知（Gnosis）\n"
        "  • 各文明秘传传统的核心（印度教的 Purusha、\n"
        "    曼达派的 Manda、希腊的 Anthropos）\n"
        "  • 目标：自我认知（Self-Knowledge）→ 与神合一\n\n"
        "世间智慧（Temporal Wisdom）：\n"
        "  • 应用高等智慧于实际生活\n"
        "  • 文明由此而生\n"
        "  • 包括科学、艺术、工艺\n\n"
        "占星学的角色：连接高等智慧与世间智慧的桥梁。\n"
        "本命盘显示一个人是否倾向于追求自我认知。"
    )),

    # Module 2: Medieval vs Modern
    ("中世纪与现代占星学的区别（一）", (
        "第一：中世纪占星学几乎没有心理学内容。\n\n"
        "为什么？\n"
        "  • 中世纪教会认为灵魂不可被星象决定\n"
        "  • 灵魂属于上帝的领域，不归行星管辖\n"
        "  • 中世纪占星著作关注的是具体事件的预测\n"
        "    —— 命主会做什么，会发生什么\n\n"
        "对比现代占星学：\n"
        "  • 大量讨论心理状态、人格发展\n"
        "  • 回避对具体事件的确定性预测\n\n"
        "Zoller 的观点：两种传统各有价值，但本课程只教传统方法。\n"
        "技术正确了，结果自然准确。"
    )),

    ("中世纪与现代占星学的区别（二）", (
        "第二：中世纪占星学理解'命运'的方式不同。\n\n"
        "中世纪人的处境（Zoller 的描述）：\n"
        "  • 人在未重生状态下受必然性（Necessity）约束\n"
        "  • 命运 = 自然法则 + 社会法则（教会、国王、阶级）\n"
        "  • 行为在相当程度上是可预测的\n\n"
        "这解释了为什么中世纪占星师相信具体预测是可能的：\n"
        "  阶级、民族、宗教都限制了一个人的可能性范围。\n"
        "  Ibn Ezra 指出：一个拥有'国王星盘'的犹太人\n"
        "  在中世纪不可能成为国王，因为他属于流散民族。\n\n"
        "现代人的假设恰恰相反：每个人都是'自由'的个体。"
    )),

    ("中世纪与现代占星学的区别（三）", (
        "第三：中世纪占星学的预测方法更加系统化。\n\n"
        "Zoller 的比喻：解盘就像翻译。\n"
        "  行星在星座中 = 主语\n"
        "  行星在宫位中 = 谓语\n"
        "  行星间的相位 = 修饰语\n\n"
        "就像语法规则一样，有一套规则将天体语言\n"
        "翻译成关于人生事件的日常语言。\n\n"
        "直觉的作用：\n"
        "  • 适当的直觉会自然进入解盘过程\n"
        "  • 但保持在最低限度 —— 过滤掉大量虚假的直觉\n"
        "  • 通过规则检验后剩下的，正确率远高于纯直觉\n\n"
        "Zoller：'我们传达正确的东西，不是因为我们是大师，\n"
        "而是因为这门艺术本身是正确的。'"
    )),

    # Module 3: Source Texts
    ("中世纪占星学的经典文献", (
        "Zoller 列举了塑造中世纪占星传统的关键文献：\n\n"
        "希腊化时期（Hellenistic）：\n"
        "  • Nechepso-Petosiris（失传，仅存残篇）\n"
        "  • 赫尔墨斯之书（Liber Hermetis）\n"
        "  • Dorotheus - Carmen Astrologicum（五经）\n"
        "  • Ptolemy - Tetrabiblos（四书）\n"
        "  • Firmicus Maternus - Mathesis\n"
        "  • Vettius Valens - Anthologies（选集）\n\n"
        "阿拉伯/波斯时期：\n"
        "  • Abu Ma'shar（阿布·马沙尔）\n"
        "  • Al-Qabisi（Alcabitius）\n"
        "  • Masha'allah\n\n"
        "拉丁中世纪：\n"
        "  • Guido Bonatti - Liber Astronomiae\n"
        "  • Placidus - Primum Mobile"
    )),

    # Module 4: Fate and Free Will
    ("命运、自由意志与占星学的目的", (
        "Zoller 的核心哲学立场：\n\n"
        "关于自由意志：\n"
        "  • '你可以成为任何你想成为的人' —— Zoller 认为这是错的\n"
        "  • 如果人人皆可如此，没有人会停留在现状\n"
        "  • 本命盘限制了你的可能性范围\n\n"
        "关于灵魂的解放：\n"
        "  • 通过自我认知（Self-Knowledge），灵魂认识到自己是神的形象\n"
        "  • 与神合一后，不再受身体的统治\n"
        "  • 人变得自由，同时活在此世和天国\n"
        "  • 但在身体寿命结束前，仍需面对命运的一切安排\n\n"
        "自由不在于改变命运，而在于改变与命运的关系。\n"
        "你不行动 —— 是'自性'（Self）在观看，\n"
        "而宇宙的占星机器在作用于你。"
    )),

    ("占星学的实用价值", (
        "在哲学基础之后，Zoller 转向占星学的实用层面：\n\n"
        "占星学能回答什么问题？\n"
        "  • 今年能赚多少钱？\n"
        "  • 什么时候结婚？\n"
        "  • 会有几个孩子？\n\n"
        "这些问题的答案带来的是 —— 心安（Peace of Mind）。\n"
        "只有失去过心安的人，才知道这有多珍贵。\n\n"
        "Zoller 发现中世纪占星学比任何其他形式的占星学\n"
        "都更准确。他乐于分享，并祝学生学习成功。\n\n"
        "导论课到此结束。接下来进入 Lesson One。"
    )),
]


def add_slide(prs, title: str, body: str):
    """Add a content slide with title and body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body_frame = slide.shapes.placeholders[1].text_frame
    body_frame.clear()

    first = True
    for line in body.strip().split("\n"):
        if first:
            p = body_frame.paragraphs[0]
            first = False
        else:
            p = body_frame.add_paragraph()
        p.text = line
        p.font.size = SMALL_FONT if line.startswith("  ") else BODY_FONT
        p.space_after = Pt(6)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Title slide ----
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = "中世纪占星学基础教程"
    ts.shapes.placeholders[1].text = (
        "Medieval Astrology Foundation Course\n"
        "Robert Zoller, 2000\n\n"
        "Orientation（导论）"
    )

    # ---- Architecture overview ----
    arch_text = (
        "Orientation 导论课的知识架构（共 4 个模块）：\n\n"
        "一、为什么需要导论？\n"
        "    中世纪占星学的赫尔墨斯语境\n"
        "    高等智慧与世间智慧的关系\n"
        "    ↓\n"
        "二、中世纪占星学与现代占星学的三大区别\n"
        "    1. 几乎不涉及心理学\n"
        "    2. 对'命运'的理解根本不同\n"
        "    3. 系统化、规则化的预测方法\n"
        "    ↓\n"
        "三、经典文献来源\n"
        "    希腊化 → 阿拉伯/波斯 → 拉丁中世纪的传承脉络\n"
        "    ↓\n"
        "四、命运、自由意志与占星学的目的\n"
        "    自由不在于改变命运，而在于改变与命运的关系"
    )
    add_slide(prs, "导论课知识架构", arch_text)

    # ---- Content slides ----
    for title, body in ORIENTATION_SLIDES:
        add_slide(prs, title, body)

    # ---- Summary ----
    summary = (
        "Orientation 导论课的核心要点：\n\n"
        "1. 中世纪占星学必须放在赫尔墨斯传统中理解\n"
        "   —— 它不是孤立的算命技术，而是通向神圣认知的途径\n\n"
        "2. 现代占星学 ≠ 中世纪占星学\n"
        "   最大的区别在于对心理学的态度和对命运的理解\n\n"
        "3. 中世纪方法更加系统化、规则化\n"
        "   解盘像翻译一样，有一套'天体语法'可以学习\n\n"
        "4. 占星学的终极目的不是预测未来\n"
        "   而是帮助人实现自我认知（Self-Knowledge）\n\n"
        "下一步：Lesson One —— 星座细分与主宰关系"
    )
    add_slide(prs, "导论课总结", summary)

    print(f"Created {len(prs.slides)} slides")
    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
